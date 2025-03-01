from pptx import Presentation
from tqdm.auto import tqdm
import shutil
from input_preprocessing.documents.DocumentPreprocessing import *
from input_preprocessing.documents.filters.extract import *


class PowerPointProcessor(DocumentProcessor):
    def extract_text_and_images(self) -> str:
        try:
            ppt = Presentation(self.path)
            self.stats["total_pages"] = len(ppt.slides)
            extracted_content = {"pages": []}

            start_time = time.time()
            
            for slide_number, slide in enumerate(tqdm(ppt.slides), start=1):
                slide_data = self._process_slide(slide, slide_number)
                if slide_data["content"]:
                    extracted_content["pages"].append(slide_data)

            filename = os.path.join(self.folder_text_path, f"{self.folder_name}.json")
            self.write_json_to_file(extracted_content, filename)
            
            self.stats["processing_time"] = time.time() - start_time
            return filename

        except Exception as e:
            print(f"Error processing PowerPoint: {e}")
            shutil.rmtree(self.folder_path, ignore_errors=True)
            return ""

    def _process_slide(self, slide, slide_number: int) -> Dict:
        slide_data = {"page_number": slide_number, "content": []}
        
        for index, shape in enumerate(slide.shapes, start=1):
            if shape.has_text_frame and shape.text_frame.text.strip():
                self._process_text_shape(shape, slide_data)
            elif hasattr(shape, "image"):
                self._process_image_shape(shape, slide_number, index, slide_data)
                
        return slide_data

    def _process_text_shape(self, shape, slide_data: Dict) -> None:
        text = shape.text_frame.text.strip()
        if text:
            slide_data["content"].append({"type": "text", "text": text})
            self.stats["total_text_blocks"] += 1

    def _process_image_shape(self, shape, slide_number: int, index: int, 
                           slide_data: Dict) -> None:
        image_data = shape.image.blob
        ocr_text = TextExtractor.extract_text_from_image(image_data)

        image_path = os.path.join(
            self.folder_image_path, f"image_{slide_number}_{index}.png"
        )
        with open(image_path, 'wb') as f:
            f.write(image_data)

        content = {"type": "image", "image_path": image_path}
        if ocr_text:
            content["ocr_text"] = ocr_text
            self.stats["total_ocr_text_blocks"] += 1
        else:
            self.stats["total_images_extracted"] += 1

        slide_data["content"].append(content)