from pptx import Presentation
from sentence_transformers import SentenceTransformer
from tqdm.auto import tqdm
import json
import time
import uuid
from extract import *
from candidate_labels import * 
import shutil

# Optional: Uncomment to download NLTK packages if not already downloaded
# nltk.download('stopwords')
# nltk.download('punkt')
# nltk.download('wordnet')

class PowerPointProcessor:
    def __init__(self, path, relevance_model):
        self.path = path
        self.relevance_model = relevance_model
        self.folder_name = path.split('/')[-1].split('.')[0] + str(uuid.uuid4())
        self.folder_path = os.path.join('Data', self.folder_name)
        self.folder_image_path = os.path.join(self.folder_path, 'Images')
        self.folder_text_path = os.path.join(self.folder_path, 'Text')
        self.stats = {
            "total_slides": 0,
            "total_images_extracted": 0,
            "total_relevant_text_blocks": 0,
            "total_ocr_text_blocks": 0,
            "processing_time": 0,
        }
        os.makedirs(self.folder_image_path, exist_ok=True)
        os.makedirs(self.folder_text_path, exist_ok=True)

    def write_json_to_file(self, data, filename):
        try:
            with open(filename, "w", encoding="utf-8") as json_file:
                json.dump(data, json_file, indent=4, ensure_ascii=False)
                print(f"JSON data has been written to {filename}")
        except Exception as e:
            print(f"An error occurred while writing JSON to file: {e}")

    def extract_text_and_images(self):
        try:
            ppt = Presentation(self.path)
            self.stats["total_slides"] = len(ppt.slides)

            extracted_content = {"slides": []}
            all_text = TextExtractor.extract_all_text(ppt)

            labels = LabelGenerator.generate_labels(all_text, self.relevance_model)
            label_embeddings = self.relevance_model.encode(labels)

            start_time = time.time()
            for slide_number, slide in enumerate(tqdm(ppt.slides, desc="Processing slides"), start=1):
                slide_data = {"slide_number": slide_number, "content": []}

                for index, shape in enumerate(slide.shapes, start=1):
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        text = shape.text_frame.text.strip()
                        preprocessed_text = TextPreprocessor.preprocess_text(text)
                        if preprocessed_text and TextPreprocessor.is_relevant_text(preprocessed_text, self.relevance_model, label_embeddings):
                            slide_data["content"].append({"type": "text", "text": text})
                            self.stats["total_relevant_text_blocks"] += 1

                    elif hasattr(shape, "image"):
                        image_data = shape.image.blob
                        ocr_text = TextExtractor.extract_text_from_image(image_data, self.relevance_model, label_embeddings)

                        image_path = os.path.join(self.folder_image_path, f"image_{slide_number}_{index}.png")
                        with open(image_path, 'wb') as f:
                            f.write(image_data)

                        content = {"type": "image", "image_path": image_path}
                        if ocr_text:
                            content["ocr_text"] = ocr_text
                            self.stats["total_ocr_text_blocks"] += 1
                        else:
                            self.stats["total_images_extracted"] += 1

                        slide_data["content"].append(content)

                if slide_data["content"]:
                    extracted_content["slides"].append(slide_data)

            filename = os.path.join(self.folder_text_path, f"Data-{self.folder_name}.json")
            self.write_json_to_file(extracted_content, filename)

            self.stats["processing_time"] = time.time() - start_time
            return filename

        except Exception as e:
            print(f"An error occurred while processing the PowerPoint: {e}")
            shutil.rmtree(self.folder_path, ignore_errors=True)

    def print_stats(self):
        print("\n=== PowerPoint Processing Stats ===")
        for key, value in self.stats.items():
            print(f"{key.replace('_', ' ').capitalize()}: {value}")




if __name__ == "__main__":
    relevance_model = SentenceTransformer('all-MiniLM-L6-v2')
    processor = PowerPointProcessor('./lec.pptx', relevance_model)
    output_file = processor.extract_text_and_images()
    processor.print_stats()
